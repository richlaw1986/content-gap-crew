"""Extract text from uploaded documents (PDF, DOCX, PPTX, TXT, CSV, MD).

Used by the knowledge-document extraction endpoint to populate the
``extractedSummary`` field on agent documents in Sanity.
"""

from __future__ import annotations

import io
import logging
from typing import Any

import httpx

logger = logging.getLogger(__name__)

# Maximum characters to keep from a single document before summarisation
MAX_EXTRACT_CHARS = 12_000


def extract_text(file_bytes: bytes, filename: str) -> str:
    """Detect file type by extension and return extracted plain text."""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    if lower.endswith(".docx"):
        return _extract_docx(file_bytes)
    if lower.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    if lower.endswith((".txt", ".md", ".csv")):
        return _extract_text(file_bytes)
    raise ValueError(f"Unsupported file type: {filename}")


# ── PDF ────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("pypdf is not installed — run: pip install pypdf")

    reader = PdfReader(io.BytesIO(data))
    pages: list[str] = []
    total = 0
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(text)
            total += len(text)
            if total > MAX_EXTRACT_CHARS:
                break
    return "\n\n".join(pages)[:MAX_EXTRACT_CHARS]


# ── DOCX ───────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document  # python-docx
    except ImportError:
        raise RuntimeError("python-docx is not installed — run: pip install python-docx")

    doc = Document(io.BytesIO(data))
    paragraphs: list[str] = []
    total = 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
            total += len(text)
            if total > MAX_EXTRACT_CHARS:
                break
    return "\n\n".join(paragraphs)[:MAX_EXTRACT_CHARS]


# ── PPTX ───────────────────────────────────────────────────────

def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        raise RuntimeError("python-pptx is not installed — run: pip install python-pptx")

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    total = 0
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slide_text = f"[Slide {i}]\n" + "\n".join(texts)
            slides.append(slide_text)
            total += len(slide_text)
            if total > MAX_EXTRACT_CHARS:
                break
    return "\n\n".join(slides)[:MAX_EXTRACT_CHARS]


# ── Plain text / Markdown / CSV ────────────────────────────────

def _extract_text(data: bytes) -> str:
    # Try UTF-8 first, fall back to latin-1
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1")
    return text.strip()[:MAX_EXTRACT_CHARS]


# ── Download from Sanity CDN ──────────────────────────────────

async def download_sanity_asset(url: str) -> bytes:
    """Download a file from Sanity's CDN."""
    async with httpx.AsyncClient(timeout=60.0) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        return resp.content


async def extract_from_sanity_asset(url: str, filename: str) -> str:
    """Download a Sanity asset and extract text from it."""
    data = await download_sanity_asset(url)
    return extract_text(data, filename)
